"""Text extraction and file-type detection helpers."""

import csv
import os
import re
import sys

from bs4 import BeautifulSoup
from markdownify import MarkdownConverter
from PyPDF2 import PdfReader
import chardet
from pptx import Presentation
from docx import Document
from lxml import etree
from openpyxl import load_workbook


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def md(soup, **options):
    return MarkdownConverter(**options).convert_soup(soup)


# ---------------------------------------------------------------------------
# Extraction functions
# ---------------------------------------------------------------------------

def extract_text_from_html(html_content):
    # Convert the modified HTML content to Markdown
    try:
        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove all <script> tags
        for script in soup.find_all('script'):
            script.decompose()

        # Remove all <style> tags
        for style in soup.find_all('style'):
            style.decompose()

        # Remove all <noscript> tags
        for noscript in soup.find_all('noscript'):
            noscript.decompose()

        # Remove all <svg> tags
        for svg in soup.find_all('svg'):
            svg.decompose()

        # Remove all <canvas> tags
        for canvas in soup.find_all('canvas'):
            canvas.decompose()
        
        # Remove all <audio> tags
        for audio in soup.find_all('audio'):
            audio.decompose()

        # Remove all <video> tags
        for video in soup.find_all('video'):
            video.decompose()

        # Remove all <iframe> tags
        for iframe in soup.find_all('iframe'):
            iframe.decompose()

        # Remove navigation elements that pollute RAG indexing
        for element in soup.find_all(['nav', 'header', 'footer', 'aside']):
            element.decompose()

        # Remove elements by common navigation-related class names
        nav_classes = ['navigation', 'nav', 'navbar', 'menu', 'sidebar',
                       'breadcrumb', 'breadcrumbs', 'toc', 'table-of-contents',
                       'site-header', 'site-footer', 'site-nav']
        for element in soup.find_all(class_=lambda c: c and any(
                nav_cls in cls.lower() for cls in (c if isinstance(c, list) else [c]) for nav_cls in nav_classes)):
            element.decompose()

        # Remove elements by common navigation-related id patterns
        nav_ids = ['nav', 'navigation', 'menu', 'sidebar', 'header', 'footer',
                   'breadcrumb', 'toc', 'table-of-contents']
        for element in soup.find_all(id=lambda i: i and any(
                nav_id in i.lower() for nav_id in nav_ids)):
            element.decompose()

        # Remove elements with common navigation ARIA roles
        for element in soup.find_all(attrs={'role': ['navigation', 'banner', 'contentinfo', 'menu', 'menubar']}):
            element.decompose()

        text = md(soup, strip=['a', 'img'], heading_style='ATX', 
                        escape_asterisks=False, escape_underscores=False, 
                        autolinks=False)
        
        # Remove extra newlines
        text = re.sub(r'\n+', '\n', text)

        return text
    except Exception as e:
        print(f"Failed to parse HTML content: {e}", file=sys.stderr)
        return ""

def extract_text_from_pdf(pdf_content):
    with open('temp.pdf', 'wb') as f:
        f.write(pdf_content)

    reader = PdfReader('temp.pdf')
    text = ''
    for page in reader.pages:
        text += page.extract_text()

    # Clean up by removing the temporary file
    os.remove('temp.pdf')

    # Return the extracted text, with extra newlines removed
    return re.sub(r'\n+', '\n', text)

def extract_text_from_docx(docx_path):
    # Load the Word document
    document = Document(docx_path)
    
    # Extract the file name (without extension) and replace underscores with spaces
    file_name = os.path.splitext(os.path.basename(docx_path))[0].replace('_', ' ')
    
    # Initialize a list to collect Markdown lines
    markdown_lines = []
    
    def process_paragraph(paragraph, list_level=0):
        """Convert a paragraph into Markdown based on its style and list level."""
        text = paragraph.text.replace("\n", " ").strip()  # Replace carriage returns with spaces
        if not text:
            return None  # Skip empty paragraphs
        
        # Check if paragraph is a list item based on indentation
        if paragraph.style.name == "List Paragraph":
            # Use the list level to determine indentation for bullet points
            bullet_prefix = "  " * list_level + "- "
            return f"{bullet_prefix}{text}"
        
        # Check for headings
        if paragraph.style.name.startswith("Heading"):
            heading_level = int(paragraph.style.name.split(" ")[1])
            return f"{'#' * heading_level} {text}"

        # Default: Regular paragraph
        return text
    
    def extract_lists(docx):
        """Extract the list structure from the underlying XML of the document."""
        # Access the document XML using lxml
        xml_content = docx.element
        namespaces = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        }

        # Parse the XML tree using lxml's etree
        root = etree.fromstring(etree.tostring(xml_content))

        # Find all list items (w:li)
        list_paragraphs = []
        for item in root.xpath("//w:li", namespaces=namespaces):
            # Extract the list level from the parent elements
            list_level = item.getparent().getparent().get("w:ilvl")
            if list_level is not None:
                list_paragraphs.append((list_level, item.text.strip()))
        
        return list_paragraphs
    
    # Add the document title (file name) as the top-level heading
    markdown_lines.append(f"# {file_name}")
    
    # Process each paragraph in the document
    for paragraph in document.paragraphs:
        # Detect bullet points based on paragraph's indent level (style `List Paragraph`)
        markdown_line = process_paragraph(paragraph)
        if markdown_line:
            markdown_lines.append(markdown_line)

    # Extract and process lists directly from the document's XML
    lists = extract_lists(document)
    for level, item in lists:
        bullet_prefix = "  " * int(level) + "- "
        markdown_lines.append(f"{bullet_prefix}{item}")
    
    # Join all lines into a single Markdown string
    return "\n\n".join(markdown_lines)

def extract_text_from_csv(csv_path):
    """Extract text from a CSV file, converting it to a Markdown table."""
    file_name = os.path.splitext(os.path.basename(csv_path))[0].replace('_', ' ')
    markdown_lines = [f"# {file_name}"]

    # Detect encoding
    with open(csv_path, 'rb') as f:
        raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get('encoding', 'utf-8') or 'utf-8'

    with open(csv_path, 'r', encoding=encoding, errors='replace') as f:
        # Sniff the dialect (delimiter, quoting, etc.)
        sample = f.read(8192)
        f.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel  # fallback to default comma-separated

        reader = csv.reader(f, dialect)
        rows = list(reader)

    if not rows:
        return f"# {file_name}\n\n(empty file)"

    # First row as table header
    headers = [cell.strip() for cell in rows[0]]
    markdown_lines.append('| ' + ' | '.join(headers) + ' |')
    markdown_lines.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')

    # Remaining rows as table data
    for row in rows[1:]:
        # Skip entirely empty rows
        if all(not cell.strip() for cell in row):
            continue
        cells = [cell.strip() for cell in row]
        # Pad or trim to match header count
        while len(cells) < len(headers):
            cells.append('')
        cells = cells[:len(headers)]
        markdown_lines.append('| ' + ' | '.join(cells) + ' |')

    return '\n'.join(markdown_lines)

def extract_text_from_xlsx(xlsx_path):
    """Extract text from an XLSX file, converting each sheet to a Markdown table."""
    workbook = load_workbook(xlsx_path, read_only=True, data_only=True)
    
    file_name = os.path.splitext(os.path.basename(xlsx_path))[0].replace('_', ' ')
    markdown_lines = [f"# {file_name}"]
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows = list(sheet.iter_rows(values_only=True))
        
        if not rows:
            continue
        
        # Add sheet name as heading
        markdown_lines.append(f"\n## {sheet_name}")
        
        # First row as table header
        headers = [str(cell) if cell is not None else '' for cell in rows[0]]
        markdown_lines.append('| ' + ' | '.join(headers) + ' |')
        markdown_lines.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')
        
        # Remaining rows as table data
        for row in rows[1:]:
            # Skip entirely empty rows
            if all(cell is None for cell in row):
                continue
            cells = [str(cell) if cell is not None else '' for cell in row]
            # Pad or trim to match header count
            while len(cells) < len(headers):
                cells.append('')
            cells = cells[:len(headers)]
            markdown_lines.append('| ' + ' | '.join(cells) + ' |')
    
    workbook.close()
    return '\n'.join(markdown_lines)

def extract_text_from_pptx(pptx_path):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_path)
    
    # Extract the file name (without extension) and replace underscores with spaces
    file_name = os.path.splitext(os.path.basename(pptx_path))[0].replace('_', ' ')
    
    # Initialize a list to collect Markdown lines
    markdown_lines = []
    
    def extract_text_with_bullets(shape, exclude_text=None):
        """Extract text with proper bullet point levels from a shape."""
        text_lines = []
        if shape.is_placeholder or shape.has_text_frame:
            if shape.text_frame and shape.text_frame.text.strip():
                for paragraph in shape.text_frame.paragraphs:
                    line_text = paragraph.text.replace("\r", "").replace("\n", " ").strip()  # Replace \n with space
                    if line_text and line_text != exclude_text:  # Exclude the slide title if needed
                        bullet_level = paragraph.level  # Get the bullet level
                        bullet = "  " * bullet_level + "- " + line_text
                        text_lines.append(bullet)
        elif shape.shape_type == 6:  # Grouped shapes
            # Handle grouped shapes recursively
            for sub_shape in shape.shapes:
                text_lines.extend(extract_text_with_bullets(sub_shape, exclude_text))
        return text_lines
    
    def get_first_text_entry(slide):
        """Retrieve the first text entry from the slide."""
        for shape in slide.shapes:
            if shape.is_placeholder or shape.has_text_frame:
                if shape.text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.paragraphs[0].text.replace("\n", " ").strip()
        return None
    
    for slide_number, slide in enumerate(presentation.slides, start=1):
        # Determine the Markdown header level
        if slide_number == 1:
            header_prefix = "#"
        else:
            header_prefix = "##"
        
        # Add the slide title or file name as the main title for the first slide
        if slide_number == 1:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
            else:
                title = file_name
            markdown_lines.append(f"{header_prefix} {title}")
        else:
            # Add the title for subsequent slides
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
            else:
                # Use the first text entry as the slide title if no title is present
                title = get_first_text_entry(slide)
                if not title:
                    title = f"Slide {slide_number}"
            markdown_lines.append(f"{header_prefix} {title}")
        
        # Add the slide content (text in other shapes), excluding the title if it's used
        for shape in slide.shapes:
            bullet_text = extract_text_with_bullets(shape, exclude_text=title)
            markdown_lines.extend(bullet_text)
        
        # Add a separator between slides, except after the last slide
        if slide_number < len(presentation.slides):
            markdown_lines.append("")
    
    # Join all lines into a single Markdown string
    return "\n".join(markdown_lines)


# ---------------------------------------------------------------------------
# File-type detection
# ---------------------------------------------------------------------------

def is_html(file_path):
    """
    Check if the given file is an HTML file, either by its extension or content.
    """
    # Check for .htm and .html extensions
    if file_path.endswith(".htm") or file_path.endswith(".html") or file_path.endswith(".xhtml"):
        return True
    
    # Check for HTML files without extensions
    try:
        with open(file_path, 'r') as f:
            first_line = next((line.strip() for line in f if line.strip()), None)
            return first_line and (first_line.lower().startswith('<!doctype html>') or first_line.lower().startswith('<html'))
    except Exception as e:
        return False
    
def is_docx(file_path):
    """
    Check if the given file is a DOCX file.
    """
    # Check for .docx extension
    if file_path.lower().endswith(".docx"):
        return True

    return False
    
def is_pptx(file_path):
    """
    Check if the given file is a PPTX file.
    """
    # Check for .pptx extension
    if file_path.lower().endswith(".pptx"):
        return True

    return False

def is_markdown(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist.")
    
    # Automatically consider .md files as Markdown
    if file_path.endswith('.md'):
        return True
    
    # If the file is not .md, but is .txt, proceed with content checking
    if not file_path.endswith('.txt'):
        return False

    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            
            # Check for common Markdown patterns
            if re.match(r'^#{1,6}\s', line):  # Heading (e.g., # Heading)
                return True
    
    # If no Markdown features are found, assume it's a regular text file
    return False
